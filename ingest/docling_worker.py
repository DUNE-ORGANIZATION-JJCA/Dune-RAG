"""
Docling Worker - Document extraction for PDF, DOCX, PPTX
"""

import os
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from dataclasses import dataclass

logger = logging.getLogger(__name__)

@dataclass
class DocumentResult:
    """Result of document processing"""
    text: str
    metadata: Dict[str, Any]
    success: bool
    error: Optional[str] = None

class DoclingWorker:
    """
    Extracts text from documents using Docling.
    Supports: PDF, DOCX, PPTX, images
    """
    
    def __init__(self):
        self.supported_formats = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg"}
    
    def extract(self, file_path: str) -> DocumentResult:
        """
        Extract text from a document.
        
        Args:
            file_path: Path to the document
            
        Returns:
            DocumentResult with extracted text and metadata
        """
        path = Path(file_path)
        
        if not path.exists():
            return DocumentResult(
                text="",
                metadata={"source": str(file_path)},
                success=False,
                error=f"File not found: {file_path}"
            )
        
        ext = path.suffix.lower()
        if ext not in self.supported_formats:
            return DocumentResult(
                text="",
                metadata={"source": str(file_path)},
                success=False,
                error=f"Unsupported format: {ext}"
            )
        
        try:
            if ext == ".pdf":
                return self._extract_pdf(path)
            elif ext == ".docx":
                return self._extract_docx(path)
            elif ext == ".pptx":
                return self._extract_pptx(path)
            else:
                return self._extract_image(path)
                
        except Exception as e:
            logger.error(f"Error extracting {file_path}: {e}")
            return DocumentResult(
                text="",
                metadata={"source": str(file_path)},
                success=False,
                error=str(e)
            )
    
    def _extract_pdf(self, path: Path) -> DocumentResult:
        """Extract text from PDF using Docling"""
        try:
            from docling.document_converter import DocumentConverter
            
            converter = DocumentConverter()
            result = converter.convert(str(path))
            
            text = result.document.export_to_text()
            
            return DocumentResult(
                text=text,
                metadata={
                    "source": str(path),
                    "doc_type": "pdf",
                    "file_name": path.name,
                    "pages": len(result.pages) if hasattr(result, 'pages') else 0
                },
                success=True
            )
        except ImportError:
            # Fallback to PyMuPDF
            return self._extract_pdf_fallback(path)
    
    def _extract_pdf_fallback(self, path: Path) -> DocumentResult:
        """Fallback PDF extraction with PyMuPDF"""
        import fitz  # PyMuPDF
        
        text_parts = []
        with fitz.open(str(path)) as doc:
            for page in doc:
                text_parts.append(page.get_text())
        
        text = "\n\n".join(text_parts)
        
        return DocumentResult(
            text=text,
            metadata={
                "source": str(path),
                "doc_type": "pdf",
                "file_name": path.name,
                "pages": len(text_parts)
            },
            success=True
        )
    
    def _extract_docx(self, path: Path) -> DocumentResult:
        """Extract text from DOCX"""
        from docx import Document
        
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        
        return DocumentResult(
            text=text,
            metadata={
                "source": str(path),
                "doc_type": "docx",
                "file_name": path.name
            },
            success=True
        )
    
    def _extract_pptx(self, path: Path) -> DocumentResult:
        """Extract text from PPTX"""
        from pptx import Presentation
        
        prs = Presentation(str(path))
        text_parts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        
        text = "\n\n".join(text_parts)
        
        return DocumentResult(
            text=text,
            metadata={
                "source": str(path),
                "doc_type": "pptx",
                "file_name": path.name,
                "slides": len(prs.slides)
            },
            success=True
        )
    
    def _extract_image(self, path: Path) -> DocumentResult:
        """OCR for images"""
        try:
            from PIL import Image
            import pytesseract
            
            image = Image.open(str(path))
            text = pytesseract.image_to_string(image, lang='spa+eng')
            
            return DocumentResult(
                text=text,
                metadata={
                    "source": str(path),
                    "doc_type": "image",
                    "file_name": path.name
                },
                success=True
            )
        except ImportError:
            return DocumentResult(
                text="",
                metadata={"source": str(path)},
                success=False,
                error="OCR not available (pytesseract not installed)"
            )
    
    def extract_batch(self, file_paths: List[str]) -> List[DocumentResult]:
        """Extract from multiple files"""
        results = []
        for path in file_paths:
            result = self.extract(path)
            results.append(result)
        return results